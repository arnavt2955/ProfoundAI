from flask import Flask, request, jsonify
from dotenv import load_dotenv
import os
import pdfplumber
from pymongo import MongoClient
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import MongoDBAtlasVectorSearch
from langchain.document_loaders import DirectoryLoader
from langchain.schema import Document
from langchain.llms import OpenAI
from langchain.chains import RetrievalQA
from pymongo import MongoClient
from canvasapi import Canvas
from canvasapi.exceptions import ResourceDoesNotExist, Unauthorized
from PyPDF2 import PdfReader
from pptx import Presentation
from io import BytesIO
import requests

load_dotenv()
OPENAI_KEY=os.getenv('OPENAI_KEY')
MONGO_URI=os.getenv('MONGO_URI')
client = MongoClient(MONGO_URI)
dbName = "profound_slides"
collectionName = "collection_of_text_blobs"
collection = client[dbName][collectionName]

embeddings = OpenAIEmbeddings(openai_api_key=OPENAI_KEY)
llm = OpenAI(openai_api_key=OPENAI_KEY, temperature=0.5)
app = Flask(__name__)

@app.route("/test")
def test():
    return {"message": "dummy"}

@app.route('/upload', methods=['POST'])
def upload_file():
    if 'file' not in request.files:
        return jsonify({"error": "No file uploaded"}), 400

    file = request.files['file']
    canvasURL = request.form.get('canvasURL')
    canvasToken = request.form.get('canvasToken')

    #file.save(f'./uploads/{file.filename}')
    
    full_text = ""
    page_text = []

    with pdfplumber.open(file) as pdf:
        for page in pdf.pages:
            # Extract text from each page
            text = page.extract_text() + "\n"
            # Optionally, extract tables (if available)
            tables = page.extract_tables()
            for table in tables:
                text += f"Table: {table}\n"
            full_text += text + "\n"
            page_text.append(text)
    # Clear the collection before adding new documents
    collection.delete_many({})  
    documents = [Document(page_content=doc) for doc in page_text]

    canvasFiles = retrieveCanvasFiles(canvasURL, canvasToken)

    if (len(canvasFiles) > 0):
        pdf_Files = [file for file in canvasFiles if file.display_name.lower().endswith('.pdf')]
        # for canvasfile in pdf_Files:
        #     print(type(canvasfile), "-----------", canvasfile)
        pdf_Downloads = download_files_to_memory(pdf_Files, canvasToken)
        pdf_Pagetext = []
        for eachPage, eachBuffer in pdf_Downloads:
            pdf_Pagetext = extract_text_from_pdf_buffer(eachPage, eachBuffer)
            pdf_Pagetext = [Document(page_content=doc) for doc in pdf_Pagetext]
            documents = documents + pdf_Pagetext

        pptx_Files = [file for file in canvasFiles if file.display_name.lower().endswith('.pptx')]
        # for canvasfile in pptx_Files:
        #     print(type(canvasfile), "-----------", canvasfile)
        pptx_Downloads = download_files_to_memory(pptx_Files, canvasToken)
        pptx_Pagetext = []
        for eachPage, eachBuffer in pptx_Downloads:
            pptx_Pagetext = extract_text_from_pptx_buffer(eachPage, eachBuffer)
            pptx_Pagetext = [Document(page_content=doc) for doc in pptx_Pagetext]
            documents = documents + pptx_Pagetext

    vectorStore = MongoDBAtlasVectorSearch.from_documents(documents, embeddings, collection=collection)
    #vectorStore = MongoDBAtlasVectorSearch(collection, embeddings)
    retriever = vectorStore.as_retriever()
    qa = RetrievalQA.from_chain_type(llm, chain_type="stuff", retriever=retriever)
    output = []
    for pt in page_text:
        output.append(qa.run("Present the following parsed lecture slide like you are a professor, concisely in paragraph format. Leave out any citations or page numbers. If the slide has little content, generate less: " + pt))
    return jsonify({"full_text": full_text, "page_text": output}), 200

# retreives all pdfs and powerpoints (pptx files)
def retrieveCanvasFiles(canvasURL, canvasAuthToken):

    if "/courses/" not in canvasURL or canvasAuthToken == "":
        return []
    
    url_parts = canvasURL.split("/courses/", 1)
    baseURL = url_parts[0]

    courseURL = url_parts[1]
    
    if "/" in courseURL:
        courseIdArr = courseURL.split("/", 1)
        courseId = int(courseIdArr[0])
    else:
        courseId = int(courseURL)
        
    
    try:
        canvas = Canvas(baseURL, canvasAuthToken)
    except Unauthorized as e:
        print("Unauthorized access to canvas. Make sure Auth Token is valid for the course")
        print(f"Details: {e}")
        return []
    except Exception as e:
        print("Unexpected error creating canvas object")
        print(f"Details: {e}")
        return []
    
    try:
        # Attempt to retrieve the course
        # courseId = 374332
        canvas = canvas.get_course(courseId)
        
        # If successful, print canvas details
        print(f"Access Token is valid and has permissions for the course.")
        print(f"Course ID: {canvas.id}")
        print(f"Course Name: {canvas.name}")
        print(f"Course Start Date: {canvas.start_at}")
        print(f"Course End Date: {canvas.end_at}")
        
        # Optionally, list a few files to further verify access
        print("\nRetrieving files in the course:")
        files = canvas.get_files()  # Limiting to first 5 files for demonstration
        # for file in files:
        #     print(f"- {file.display_name} (ID: {file.id})")

        # filter to only allow powerpoint and pdfs
        allowed_extensions = ('.pdf', '.pptx')
        filtered = [file for file in files if file.display_name.lower().endswith(allowed_extensions)]

        return filtered
        
    except Unauthorized as e:
        print("Invalid token for course")
        print(f"Details: {e}")
        return []
    except Exception as e:
        print("Unexpected error retreiving course")
        print(f"Details: {e}")
        return []
    
def download_files_to_memory(files, canvasToken):
    """
    Downloads files from Canvas into in-memory buffers.
    
    Args:
        files (list): List of file objects to download.
    
    Returns:
        list: List of tuples containing file name and BytesIO buffer.
    """
    downloaded_files = []
    headers = {
        'Authorization': f'Bearer {canvasToken}'
    }
    for file in files:
        try:
            # Create a BytesIO buffer
            file_buffer = BytesIO()

            response = requests.get(file.url, headers, stream=True)
            if response.status_code == 200:
            # Write the content to the buffer in chunks
                for chunk in response.iter_content(chunk_size=8192):
                    if chunk:
                        file_buffer.write(chunk)
            else:
                print("get download response failed")
                raise Exception(response.content)
                
            
            # # Download the file content into the buffer
            # file.download(file_buffer)
            
            # Seek to the beginning of the buffer
            file_buffer.seek(0)
            
            # Append to the list as (file_name, buffer)
            downloaded_files.append((file.display_name, file_buffer))
            print(f"Successfully downloaded {file.display_name} into memory.\n")
        except Exception as e:
            print(f"Failed to download {file.display_name}: {e}\n")
    return downloaded_files

def extract_text_from_pdf_buffer(file_name, file_buffer):
    """
    Extracts text from each page of a PDF file stored in a BytesIO buffer.
    
    Args:
        file_name (str): Name of the PDF file.
        file_buffer (BytesIO): In-memory buffer containing PDF data.
    
    Returns:
        list: A list of strings, each representing text from a page.
    """
    text_pages = []
    try:
        reader = PdfReader(file_buffer)
        num_pages = len(reader.pages)
        for page_num in range(num_pages):
            page = reader.pages[page_num]
            text = page.extract_text()
            if text:
                text_pages.append(text)
            else:
                text_pages.append("")  # Empty string if no text found
    except Exception as e:
        print(f"Error extracting text from {file_name}: {e}")
    return text_pages

def extract_text_from_pptx_buffer(file_name, file_buffer):
    """
    Extracts text from each slide of a PowerPoint file stored in a BytesIO buffer.
    
    Args:
        file_name (str): Name of the PowerPoint file.
        file_buffer (BytesIO): In-memory buffer containing PowerPoint data.
    
    Returns:
        list: A list of strings, each representing text from a slide.
    """
    text_slides = []
    try:
        prs = Presentation(file_buffer)
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = ""
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text += shape.text + "\n"
            text_slides.append(slide_text.strip())
    except Exception as e:
        print(f"Error extracting text from {file_name}: {e}")
    return text_slides

@app.route('/question', methods=['POST'])
def question():
    data = request.get_json()
    message = data.get('message')
    vectorStore = MongoDBAtlasVectorSearch(collection, embeddings)
    retriever = vectorStore.as_retriever()
    qa = RetrievalQA.from_chain_type(llm, chain_type="stuff", retriever=retriever)
    out = qa.run(message)
    return jsonify({"answer": out}), 200

if __name__ == "__main__":
    print('hello')
    app.run(debug=True)
